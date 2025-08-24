"""
Multimodal document processing pipeline.
Handles file type detection, routing, and content extraction from various formats.
"""

import asyncio
import mimetypes
import tempfile
import os
import io
from pathlib import Path
from typing import Dict, List, Optional, Any, BinaryIO, Union, Tuple
from datetime import datetime
import hashlib

# Document processing libraries
from PIL import Image
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

from ..config import get_settings
from ..utils.logging import LoggingMixin, log_function_call
from ..core.exceptions import FileProcessingError, UnsupportedFileTypeError
from ..models.schemas import ContentType, ContentChunk, ProcessingStatus
from .azure_ai_service import AzureAIService


class DocumentProcessor(LoggingMixin):
    """Multimodal document processing service."""
    
    # Supported file types and their MIME types
    SUPPORTED_TYPES = {
        # Text documents
        'text/plain': ContentType.TEXT,
        'application/pdf': ContentType.DOCUMENT,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ContentType.DOCUMENT,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': ContentType.DOCUMENT,
        'application/vnd.ms-excel': ContentType.DOCUMENT,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ContentType.DOCUMENT,
        
        # Images
        'image/jpeg': ContentType.IMAGE,
        'image/jpg': ContentType.IMAGE,
        'image/png': ContentType.IMAGE,
        'image/bmp': ContentType.IMAGE,
        'image/tiff': ContentType.IMAGE,
        'image/gif': ContentType.IMAGE,
        'image/webp': ContentType.IMAGE,
        
        # Audio/Video (basic support)
        'audio/wav': ContentType.TEXT,  # Will be processed as text after transcription
        'audio/mp3': ContentType.TEXT,
        'audio/m4a': ContentType.TEXT,
        'video/mp4': ContentType.TEXT,
        'video/avi': ContentType.TEXT,
    }
    
    def __init__(self):
        """Initialize the document processor."""
        self.settings = get_settings()
        self.azure_ai = AzureAIService()
        self.processing_settings = self.settings.processing
        
        # Text chunking parameters
        self.chunk_size = getattr(self.processing_settings, 'chunk_size', 1000)
        self.chunk_overlap = getattr(self.processing_settings, 'chunk_overlap', 200)
        self.max_file_size = getattr(self.processing_settings, 'max_file_size_mb', 50) * 1024 * 1024  # Convert to bytes
    
    @log_function_call
    async def process_file(
        self,
        file_data: bytes,
        filename: str,
        source_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[ContentChunk]:
        """
        Process a file and extract content chunks.
        
        Args:
            file_data: Raw file data
            filename: Original filename
            source_id: Unique identifier for the source
            metadata: Additional metadata
            
        Returns:
            List of processed content chunks
        """
        if len(file_data) > self.max_file_size:
            raise FileProcessingError(
                f"File size ({len(file_data)} bytes) exceeds maximum allowed size ({self.max_file_size} bytes)",
                context={"filename": filename, "file_size": len(file_data)}
            )
        
        # Generate source ID if not provided
        if not source_id:
            source_id = self._generate_source_id(file_data, filename)
        
        # Initialize metadata
        if metadata is None:
            metadata = {}
        
        metadata.update({
            "filename": filename,
            "file_size": len(file_data),
            "processed_at": datetime.utcnow().isoformat(),
            "source_id": source_id
        })
        
        # Detect content type
        content_type, mime_type = self._detect_content_type(filename, file_data)
        metadata["mime_type"] = mime_type
        metadata["detected_content_type"] = content_type.value
        
        self.logger.info(
            "Processing file",
            filename=filename,
            content_type=content_type.value,
            mime_type=mime_type,
            file_size=len(file_data)
        )
        
        try:
            # Route to appropriate processor
            if content_type == ContentType.TEXT:
                if mime_type.startswith('audio/') or mime_type.startswith('video/'):
                    chunks = await self._process_audio_video(file_data, filename, source_id, metadata)
                else:
                    chunks = await self._process_text_document(file_data, filename, source_id, metadata)
            elif content_type == ContentType.DOCUMENT:
                chunks = await self._process_document(file_data, filename, source_id, metadata, mime_type)
            elif content_type == ContentType.IMAGE:
                chunks = await self._process_image(file_data, filename, source_id, metadata)
            else:
                raise UnsupportedFileTypeError(
                    f"Unsupported content type: {content_type.value}",
                    context={"filename": filename, "mime_type": mime_type}
                )
            
            self.logger.info(
                "File processing completed",
                filename=filename,
                chunks_created=len(chunks),
                content_type=content_type.value
            )
            
            return chunks
            
        except Exception as e:
            self.logger.error(
                "File processing failed",
                filename=filename,
                content_type=content_type.value,
                error=str(e)
            )
            raise FileProcessingError(
                f"Failed to process file {filename}: {str(e)}",
                context={
                    "filename": filename,
                    "content_type": content_type.value,
                    "mime_type": mime_type,
                    "original_error": str(e)
                }
            )
    
    def _detect_content_type(self, filename: str, file_data: bytes) -> Tuple[ContentType, str]:
        """
        Detect content type from filename and file data.
        
        Args:
            filename: Original filename
            file_data: Raw file data
            
        Returns:
            Tuple of (ContentType, mime_type)
        """
        # First try to detect from filename
        mime_type, _ = mimetypes.guess_type(filename)
        
        if not mime_type:
            # Try to detect from file signature
            mime_type = self._detect_mime_from_signature(file_data)
        
        if not mime_type:
            # Default to plain text
            mime_type = 'text/plain'
        
        # Map MIME type to ContentType
        content_type = self.SUPPORTED_TYPES.get(mime_type)
        
        if not content_type:
            # Check for similar types
            if mime_type.startswith('text/'):
                content_type = ContentType.TEXT
            elif mime_type.startswith('image/'):
                content_type = ContentType.IMAGE
            elif mime_type in ['application/msword', 'application/rtf']:
                content_type = ContentType.DOCUMENT
            else:
                raise UnsupportedFileTypeError(
                    f"Unsupported MIME type: {mime_type}",
                    context={"filename": filename, "mime_type": mime_type}
                )
        
        return content_type, mime_type
    
    def _detect_mime_from_signature(self, file_data: bytes) -> Optional[str]:
        """Detect MIME type from file signature."""
        if not file_data:
            return None
        
        # Common file signatures
        signatures = {
            b'\x89PNG\r\n\x1a\n': 'image/png',
            b'\xff\xd8\xff': 'image/jpeg',
            b'GIF8': 'image/gif',
            b'BM': 'image/bmp',
            b'%PDF': 'application/pdf',
            b'PK\x03\x04': 'application/zip',  # Could be DOCX, XLSX, PPTX
            b'\xd0\xcf\x11\xe0': 'application/vnd.ms-office',  # Old Office formats
        }
        
        for signature, mime_type in signatures.items():
            if file_data.startswith(signature):
                # Special handling for ZIP-based formats
                if mime_type == 'application/zip' and len(file_data) > 100:
                    # Check for Office formats
                    if b'word/' in file_data[:1000]:
                        return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                    elif b'xl/' in file_data[:1000]:
                        return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                    elif b'ppt/' in file_data[:1000]:
                        return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                return mime_type
        
        return None
    
    def _generate_source_id(self, file_data: bytes, filename: str) -> str:
        """Generate a unique source ID for the file."""
        # Create hash from file content and metadata
        hasher = hashlib.sha256()
        hasher.update(file_data)
        hasher.update(filename.encode('utf-8'))
        hasher.update(datetime.utcnow().isoformat().encode('utf-8'))
        return hasher.hexdigest()[:16]
    
    async def _process_text_document(
        self,
        file_data: bytes,
        filename: str,
        source_id: str,
        metadata: Dict[str, Any]
    ) -> List[ContentChunk]:
        """Process plain text documents."""
        try:
            # Decode text content
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            text_content = None
            
            for encoding in encodings:
                try:
                    text_content = file_data.decode(encoding)
                    metadata["encoding"] = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                raise FileProcessingError(
                    "Unable to decode text file with any supported encoding",
                    context={"filename": filename, "tried_encodings": encodings}
                )
            
            # Create chunks from text
            chunks = self._create_text_chunks(text_content, source_id, metadata)
            
            # Generate embeddings for each chunk
            for chunk in chunks:
                chunk.embedding = await self.azure_ai.generate_text_embedding(chunk.content)
            
            return chunks
            
        except Exception as e:
            raise FileProcessingError(
                f"Failed to process text document: {str(e)}",
                context={"filename": filename, "source_id": source_id}
            )
    
    async def _process_document(
        self,
        file_data: bytes,
        filename: str,
        source_id: str,
        metadata: Dict[str, Any],
        mime_type: str
    ) -> List[ContentChunk]:
        """Process structured documents (PDF, DOCX, etc.)."""
        text_content = ""
        additional_metadata = {}
        
        try:
            if mime_type == 'application/pdf':
                text_content, additional_metadata = await self._extract_pdf_content(file_data)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                text_content, additional_metadata = await self._extract_docx_content(file_data)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                text_content, additional_metadata = await self._extract_pptx_content(file_data)
            elif mime_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                text_content, additional_metadata = await self._extract_excel_content(file_data)
            else:
                # Use Azure Document Intelligence for unsupported formats
                doc_analysis = await self.azure_ai.analyze_document(file_data)
                text_content = doc_analysis.get("content", "")
                additional_metadata = {
                    "azure_analysis": doc_analysis,
                    "extraction_method": "azure_document_intelligence"
                }
            
            # Merge additional metadata
            metadata.update(additional_metadata)
            
            # Create chunks from extracted text
            chunks = self._create_text_chunks(text_content, source_id, metadata)
            
            # Generate embeddings for each chunk
            for chunk in chunks:
                chunk.embedding = await self.azure_ai.generate_text_embedding(chunk.content)
            
            return chunks
            
        except Exception as e:
            raise FileProcessingError(
                f"Failed to process document: {str(e)}",
                context={"filename": filename, "mime_type": mime_type, "source_id": source_id}
            )
    
    async def _process_image(
        self,
        file_data: bytes,
        filename: str,
        source_id: str,
        metadata: Dict[str, Any]
    ) -> List[ContentChunk]:
        """Process image files."""
        try:
            # Analyze image with Azure Computer Vision
            vision_analysis = await self.azure_ai.analyze_image(file_data)
            ocr_results = await self.azure_ai.extract_text_from_image(file_data)
            
            # Extract image metadata
            try:
                with Image.open(io.BytesIO(file_data)) as img:
                    metadata.update({
                        "image_width": img.width,
                        "image_height": img.height,
                        "image_format": img.format,
                        "image_mode": img.mode
                    })
            except Exception as e:
                self.logger.warning("Failed to extract image metadata", error=str(e))
            
            # Create content from image analysis
            content_parts = []
            
            # Add description
            descriptions = vision_analysis.get("descriptions", [])
            if descriptions:
                content_parts.append(f"Image description: {descriptions[0].get('text', '')}")
            
            # Add tags
            tags = vision_analysis.get("tags", [])
            if tags:
                tag_list = [tag["name"] for tag in tags if tag.get("confidence", 0) > 0.5]
                content_parts.append(f"Image tags: {', '.join(tag_list)}")
            
            # Add OCR text
            ocr_text = ocr_results.get("text", "").strip()
            if ocr_text:
                content_parts.append(f"Text in image: {ocr_text}")
            
            # Add object detection results
            objects = vision_analysis.get("objects", [])
            if objects:
                object_list = [obj["object"] for obj in objects if obj.get("confidence", 0) > 0.5]
                content_parts.append(f"Objects detected: {', '.join(set(object_list))}")
            
            content_text = " ".join(content_parts) if content_parts else "Image with no detectable content"
            
            # Store analysis results in metadata
            metadata.update({
                "vision_analysis": vision_analysis,
                "ocr_results": ocr_results,
                "extraction_method": "azure_computer_vision"
            })
            
            # Create a single chunk for the image
            chunk = ContentChunk(
                id=f"{source_id}_0",
                content=content_text,
                content_type=ContentType.IMAGE,
                metadata=metadata,
                source_id=source_id,
                chunk_index=0
            )
            
            # Generate embedding for the content
            chunk.embedding = await self.azure_ai.generate_text_embedding(content_text)
            
            return [chunk]
            
        except Exception as e:
            raise FileProcessingError(
                f"Failed to process image: {str(e)}",
                context={"filename": filename, "source_id": source_id}
            )
    
    async def _process_audio_video(
        self,
        file_data: bytes,
        filename: str,
        source_id: str,
        metadata: Dict[str, Any]
    ) -> List[ContentChunk]:
        """Process audio/video files (basic transcription support)."""
        try:
            # Save to temporary file for speech services
            with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name
            
            try:
                # Transcribe audio using Azure Speech Services
                transcribed_text = await self.azure_ai.speech_to_text(temp_file_path)
                
                metadata.update({
                    "transcription_service": "azure_speech_services",
                    "original_media_type": "audio/video",
                    "extraction_method": "speech_to_text"
                })
                
                # Create chunks from transcribed text
                chunks = self._create_text_chunks(transcribed_text, source_id, metadata)
                
                # Generate embeddings for each chunk
                for chunk in chunks:
                    chunk.embedding = await self.azure_ai.generate_text_embedding(chunk.content)
                
                return chunks
                
            finally:
                # Clean up temporary file
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
            
        except Exception as e:
            raise FileProcessingError(
                f"Failed to process audio/video: {str(e)}",
                context={"filename": filename, "source_id": source_id}
            )
    
    async def _extract_pdf_content(self, file_data: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PDF files."""
        try:
            import io
            
            with io.BytesIO(file_data) as pdf_stream:
                reader = pypdf.PdfReader(pdf_stream)
                
                text_content = ""
                metadata = {
                    "page_count": len(reader.pages),
                    "extraction_method": "pypdf"
                }
                
                # Extract text from each page
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content += f"\\n[Page {page_num + 1}]\\n{page_text}\\n"
                
                # Extract PDF metadata
                if reader.metadata:
                    pdf_metadata = {}
                    for key, value in reader.metadata.items():
                        if value:
                            pdf_metadata[key.replace('/', '')] = str(value)
                    metadata["pdf_metadata"] = pdf_metadata
                
                return text_content.strip(), metadata
                
        except Exception as e:
            # Fallback to Azure Document Intelligence
            self.logger.warning("PyPDF extraction failed, using Azure Document Intelligence", error=str(e))
            doc_analysis = await self.azure_ai.analyze_document(file_data)
            return doc_analysis.get("content", ""), {
                "extraction_method": "azure_document_intelligence_fallback",
                "azure_analysis": doc_analysis
            }
    
    async def _extract_docx_content(self, file_data: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract content from DOCX files."""
        try:
            import io
            
            with io.BytesIO(file_data) as docx_stream:
                doc = DocxDocument(docx_stream)
                
                text_content = ""
                metadata = {
                    "paragraph_count": len(doc.paragraphs),
                    "extraction_method": "python_docx"
                }
                
                # Extract text from paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_content += para.text + "\\n"
                
                # Extract text from tables
                table_texts = []
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            table_texts.append(row_text)
                
                if table_texts:
                    text_content += "\\n[Tables]\\n" + "\\n".join(table_texts)
                    metadata["table_count"] = len(doc.tables)
                
                return text_content.strip(), metadata
                
        except Exception as e:
            # Fallback to Azure Document Intelligence
            self.logger.warning("DOCX extraction failed, using Azure Document Intelligence", error=str(e))
            doc_analysis = await self.azure_ai.analyze_document(file_data)
            return doc_analysis.get("content", ""), {
                "extraction_method": "azure_document_intelligence_fallback",
                "azure_analysis": doc_analysis
            }
    
    async def _extract_pptx_content(self, file_data: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PowerPoint files."""
        try:
            import io
            
            with io.BytesIO(file_data) as pptx_stream:
                prs = Presentation(pptx_stream)
                
                text_content = ""
                metadata = {
                    "slide_count": len(prs.slides),
                    "extraction_method": "python_pptx"
                }
                
                # Extract text from slides
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = f"\\n[Slide {slide_num + 1}]\\n"
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text + "\\n"
                    
                    if slide_text.strip() != f"\\n[Slide {slide_num + 1}]\\n":
                        text_content += slide_text
                
                return text_content.strip(), metadata
                
        except Exception as e:
            # Fallback to Azure Document Intelligence
            self.logger.warning("PPTX extraction failed, using Azure Document Intelligence", error=str(e))
            doc_analysis = await self.azure_ai.analyze_document(file_data)
            return doc_analysis.get("content", ""), {
                "extraction_method": "azure_document_intelligence_fallback",
                "azure_analysis": doc_analysis
            }
    
    async def _extract_excel_content(self, file_data: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract content from Excel files."""
        try:
            import io
            
            with io.BytesIO(file_data) as excel_stream:
                # Try to read as Excel file
                try:
                    df_dict = pd.read_excel(excel_stream, sheet_name=None)
                except Exception:
                    # Try as CSV
                    excel_stream.seek(0)
                    df_dict = {"Sheet1": pd.read_csv(excel_stream)}
                
                text_content = ""
                metadata = {
                    "sheet_count": len(df_dict),
                    "extraction_method": "pandas"
                }
                
                # Convert each sheet to text
                for sheet_name, df in df_dict.items():
                    if not df.empty:
                        text_content += f"\\n[Sheet: {sheet_name}]\\n"
                        # Convert to string representation
                        text_content += df.to_string(index=False) + "\\n"
                
                metadata["total_rows"] = sum(len(df) for df in df_dict.values())
                metadata["total_columns"] = sum(len(df.columns) for df in df_dict.values())
                
                return text_content.strip(), metadata
                
        except Exception as e:
            # Fallback to Azure Document Intelligence
            self.logger.warning("Excel extraction failed, using Azure Document Intelligence", error=str(e))
            doc_analysis = await self.azure_ai.analyze_document(file_data)
            return doc_analysis.get("content", ""), {
                "extraction_method": "azure_document_intelligence_fallback",
                "azure_analysis": doc_analysis
            }
    
    def _create_text_chunks(
        self,
        text: str,
        source_id: str,
        metadata: Dict[str, Any]
    ) -> List[ContentChunk]:
        """Create text chunks from long text content."""
        if not text.strip():
            return []
        
        chunks = []
        text_length = len(text)
        
        # If text is short enough, create a single chunk
        if text_length <= self.chunk_size:
            chunk = ContentChunk(
                id=f"{source_id}_0",
                content=text.strip(),
                content_type=metadata.get("detected_content_type", ContentType.TEXT),
                metadata=metadata.copy(),
                source_id=source_id,
                chunk_index=0
            )
            chunks.append(chunk)
            return chunks
        
        # Split into chunks with overlap
        chunk_index = 0
        start = 0
        
        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            
            # Try to break at sentence or paragraph boundaries
            if end < text_length:
                # Look for sentence endings
                for break_char in ['.\\n', '!\\n', '?\\n', '\\n\\n']:
                    break_pos = text.rfind(break_char, start, end)
                    if break_pos > start + self.chunk_size // 2:  # Don't break too early
                        end = break_pos + len(break_char)
                        break
                
                # If no good break found, look for any whitespace
                if end == start + self.chunk_size:
                    break_pos = text.rfind(' ', start, end)
                    if break_pos > start + self.chunk_size // 2:
                        end = break_pos
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:  # Only create non-empty chunks
                chunk_metadata = metadata.copy()
                chunk_metadata.update({
                    "chunk_start": start,
                    "chunk_end": end,
                    "total_chunks": None,  # Will be updated after all chunks are created
                    "text_length": len(chunk_text)
                })
                
                chunk = ContentChunk(
                    id=f"{source_id}_{chunk_index}",
                    content=chunk_text,
                    content_type=metadata.get("detected_content_type", ContentType.TEXT),
                    metadata=chunk_metadata,
                    source_id=source_id,
                    chunk_index=chunk_index
                )
                chunks.append(chunk)
                chunk_index += 1
            
            # Move start position with overlap
            start = max(end - self.chunk_overlap, start + 1)
            
            # Prevent infinite loop
            if start >= end:
                start = end
        
        # Update total chunks count in metadata
        for chunk in chunks:
            chunk.metadata["total_chunks"] = len(chunks)
        
        return chunks
    
    @log_function_call
    async def get_supported_types(self) -> Dict[str, List[str]]:
        """
        Get list of supported file types.
        
        Returns:
            Dictionary mapping content types to supported MIME types
        """
        supported = {}
        for mime_type, content_type in self.SUPPORTED_TYPES.items():
            if content_type.value not in supported:
                supported[content_type.value] = []
            supported[content_type.value].append(mime_type)
        
        return supported
    
    @log_function_call
    async def validate_file(self, filename: str, file_size: int) -> Dict[str, Any]:
        """
        Validate if file can be processed.
        
        Args:
            filename: Original filename
            file_size: File size in bytes
            
        Returns:
            Validation result
        """
        validation_result = {
            "valid": True,
            "errors": [],
            "warnings": []
        }
        
        # Check file size
        if file_size > self.max_file_size:
            validation_result["valid"] = False
            validation_result["errors"].append(
                f"File size ({file_size} bytes) exceeds maximum allowed size ({self.max_file_size} bytes)"
            )
        
        # Check file type
        mime_type, _ = mimetypes.guess_type(filename)
        if mime_type and mime_type not in self.SUPPORTED_TYPES:
            validation_result["valid"] = False
            validation_result["errors"].append(f"Unsupported file type: {mime_type}")
        elif not mime_type:
            validation_result["warnings"].append("Unable to determine file type from filename")
        
        return validation_result